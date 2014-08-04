from pptx.shapes.table import Table
from pptx.shapes.table import _Row

class PresentationParser(object):
  
  def __init__(self, contentHandler):
    self.contentHandler = contentHandler

  def parse(self, presentation):
    for slide in presentation.slides:
        for shape in slide.shapes:
            self.__parseShape(shape)

  def __parseShape(self, shape):
    self.contentHandler.shape(shape)
    if isinstance(shape, Table):
      self.__parseTable(shape)
    if shape.has_textframe:
      self.__parseTextframe(shape.textframe)

  def __parseTable(self, table):
    self.contentHandler.table(table)
    for row in table.rows:
      self.contentHandler.tableRow(row)
      for cell in row.cells:
        self.contentHandler.tableCell(cell)
        self.__parseTextframe(cell.textframe)

  def __parseTextframe(self, textframe):
    for paragraph in textframe.paragraphs:
      self.contentHandler.paragraph(paragraph, reduce(lambda memo,run: memo + run.text, paragraph.runs, ''))

